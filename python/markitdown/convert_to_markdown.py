from pptx import Presentation
import os

def convert_all_pptx_to_markdown():
    current_dir = os.getcwd()
    files = [f for f in os.listdir(current_dir) if os.path.isfile(f) and f.endswith('.pptx')]

    if not files:
        print("変換可能なPowerPointファイルがありません。")
        return

    for pptx_file in files:
        output_file = f"{os.path.splitext(pptx_file)[0]}.md"
        try:
            prs = Presentation(pptx_file)
            with open(output_file, "w", encoding="utf-8") as md_file:
                for i, slide in enumerate(prs.slides):
                    md_file.write(f"<!-- Slide number: {i + 1} -->\n")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text = shape.text.strip()
                            md_file.write(f"{text}\n\n")
            print(f"変換成功: {pptx_file} -> {output_file}")
        except Exception as e:
            print(f"変換失敗: {pptx_file} (エラー: {e})")

if __name__ == "__main__":
    convert_all_pptx_to_markdown()
